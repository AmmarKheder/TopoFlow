#!/usr/bin/env python3
"""
Create PowerPoint Presentation for TopoFlow Project
===================================================

Generates a professional PPTX with:
- Methodology slide with diagrams
- Results slide with performance charts
- China maps with Ground Truth vs Predictions

Author: Ammar Kheddar
Project: TopoFlow
"""

import json
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import io

# Load metrics
METRICS_FILE = "archive/results_old/eval_baseline_20250923_024726/baseline_metrics.json"
MEDIA_DIR = Path("archive/media")

def load_metrics():
    """Load evaluation metrics from JSON."""
    with open(METRICS_FILE, 'r') as f:
        return json.load(f)

def create_rmse_chart(metrics):
    """Create RMSE evolution chart for all pollutants."""
    fig, ax = plt.subplots(figsize=(12, 7))

    pollutants = ['pm25', 'pm10', 'so2', 'no2', 'co', 'o3']
    pollutant_labels = ['PM2.5', 'PM10', 'SO₂', 'NO₂', 'CO', 'O₃']
    colors = ['#e74c3c', '#f39c12', '#3498db', '#9b59b6', '#1abc9c', '#2ecc71']
    horizons = ['12', '24', '48', '96']
    horizon_labels = ['12h', '24h', '48h', '96h']

    x = np.arange(len(horizons))
    width = 0.14

    for i, (pol, label, color) in enumerate(zip(pollutants, pollutant_labels, colors)):
        rmse_values = [metrics[pol][h]['rmse'] for h in horizons]
        offset = (i - 2.5) * width
        ax.bar(x + offset, rmse_values, width, label=label, color=color, alpha=0.8)

    ax.set_xlabel('Forecast Horizon', fontsize=14, fontweight='bold')
    ax.set_ylabel('RMSE (µg/m³)', fontsize=14, fontweight='bold')
    ax.set_title('TopoFlow Performance: RMSE by Pollutant and Forecast Horizon\n(Test Year 2018 - China Region)',
                 fontsize=16, fontweight='bold', pad=20)
    ax.set_xticks(x)
    ax.set_xticklabels(horizon_labels, fontsize=12)
    ax.legend(loc='upper left', fontsize=11, framealpha=0.9)
    ax.grid(axis='y', alpha=0.3, linestyle='--')
    ax.set_axisbelow(True)

    plt.tight_layout()

    # Save to BytesIO
    img_stream = io.BytesIO()
    plt.savefig(img_stream, format='png', dpi=300, bbox_inches='tight')
    plt.close()
    img_stream.seek(0)

    return img_stream

def create_mae_chart(metrics):
    """Create MAE evolution chart for selected pollutants."""
    fig, axes = plt.subplots(2, 3, figsize=(15, 10))
    fig.suptitle('TopoFlow: Mean Absolute Error (MAE) Evolution\n(Test Year 2018)',
                 fontsize=18, fontweight='bold', y=0.98)

    pollutants = ['pm25', 'pm10', 'so2', 'no2', 'co', 'o3']
    pollutant_labels = ['PM2.5', 'PM10', 'SO₂', 'NO₂', 'CO', 'O₃']
    colors = ['#e74c3c', '#f39c12', '#3498db', '#9b59b6', '#1abc9c', '#2ecc71']
    horizons = ['12', '24', '48', '96']
    horizon_x = [12, 24, 48, 96]

    for idx, (pol, label, color) in enumerate(zip(pollutants, pollutant_labels, colors)):
        ax = axes[idx // 3, idx % 3]

        mae_values = [metrics[pol][h]['mae'] for h in horizons]

        ax.plot(horizon_x, mae_values, marker='o', linewidth=3, markersize=10,
                color=color, label=label)
        ax.fill_between(horizon_x, mae_values, alpha=0.2, color=color)

        ax.set_xlabel('Forecast Horizon (hours)', fontsize=11, fontweight='bold')
        ax.set_ylabel('MAE (µg/m³)', fontsize=11, fontweight='bold')
        ax.set_title(label, fontsize=13, fontweight='bold', color=color)
        ax.grid(True, alpha=0.3, linestyle='--')
        ax.set_xlim([8, 100])

        # Add value labels
        for x, y in zip(horizon_x, mae_values):
            ax.text(x, y, f'{y:.2f}', ha='center', va='bottom', fontsize=9, fontweight='bold')

    plt.tight_layout()

    img_stream = io.BytesIO()
    plt.savefig(img_stream, format='png', dpi=300, bbox_inches='tight')
    plt.close()
    img_stream.seek(0)

    return img_stream

def create_architecture_diagram():
    """Create model architecture flow diagram."""
    fig, ax = plt.subplots(figsize=(14, 6))

    # Architecture blocks
    blocks = [
        {'name': 'Input Data\n15 variables\n128×256', 'color': '#3498db'},
        {'name': 'Wind-Guided\nScanning\n16 sectors', 'color': '#9b59b6'},
        {'name': 'Block 0\nPhysics\nAttention', 'color': '#e74c3c'},
        {'name': 'Blocks 1-5\nStandard\nTransformer', 'color': '#1abc9c'},
        {'name': 'Decoder\n6 pollutants\n4 horizons', 'color': '#f39c12'}
    ]

    block_width = 2.5
    block_height = 1.5
    spacing = 0.5
    y_center = 2

    for i, block in enumerate(blocks):
        x = i * (block_width + spacing)

        # Draw block
        rect = mpatches.FancyBboxPatch((x, y_center - block_height/2), block_width, block_height,
                                       boxstyle="round,pad=0.1",
                                       facecolor=block['color'],
                                       edgecolor='white',
                                       linewidth=3,
                                       alpha=0.9)
        ax.add_patch(rect)

        # Add text
        ax.text(x + block_width/2, y_center, block['name'],
                ha='center', va='center', fontsize=11, fontweight='bold',
                color='white', multialignment='center')

        # Draw arrow
        if i < len(blocks) - 1:
            arrow_x = x + block_width + 0.1
            arrow_end = x + block_width + spacing - 0.1
            ax.annotate('', xy=(arrow_end, y_center), xytext=(arrow_x, y_center),
                       arrowprops=dict(arrowstyle='->', lw=4, color='#2c3e50'))

    ax.set_xlim(-0.5, len(blocks) * (block_width + spacing) - spacing + 0.5)
    ax.set_ylim(0, 4)
    ax.axis('off')
    ax.set_title('TopoFlow Architecture Pipeline', fontsize=18, fontweight='bold', pad=20)

    plt.tight_layout()

    img_stream = io.BytesIO()
    plt.savefig(img_stream, format='png', dpi=300, bbox_inches='tight', facecolor='white')
    plt.close()
    img_stream.seek(0)

    return img_stream

def create_comparison_table(metrics):
    """Create performance comparison table."""
    fig, ax = plt.subplots(figsize=(14, 8))
    ax.axis('tight')
    ax.axis('off')

    pollutants = ['PM2.5', 'PM10', 'SO₂', 'NO₂', 'CO', 'O₃']
    horizons = ['12h', '24h', '48h', '96h']

    # Prepare data
    table_data = [['Pollutant'] + horizons]

    pol_keys = ['pm25', 'pm10', 'so2', 'no2', 'co', 'o3']
    for pol_label, pol_key in zip(pollutants, pol_keys):
        row = [pol_label]
        for h in ['12', '24', '48', '96']:
            rmse = metrics[pol_key][h]['rmse']
            mae = metrics[pol_key][h]['mae']
            row.append(f'{rmse:.2f}\n({mae:.2f})')
        table_data.append(row)

    table = ax.table(cellText=table_data, cellLoc='center', loc='center',
                     bbox=[0, 0, 1, 1])

    table.auto_set_font_size(False)
    table.set_fontsize(11)
    table.scale(1, 2.5)

    # Style header
    for i in range(5):
        cell = table[(0, i)]
        cell.set_facecolor('#2c3e50')
        cell.set_text_props(weight='bold', color='white', fontsize=13)

    # Style pollutant names
    colors = ['#e74c3c', '#f39c12', '#3498db', '#9b59b6', '#1abc9c', '#2ecc71']
    for i, color in enumerate(colors, start=1):
        cell = table[(i, 0)]
        cell.set_facecolor(color)
        cell.set_text_props(weight='bold', color='white', fontsize=12)

    # Style data cells
    for i in range(1, 7):
        for j in range(1, 5):
            cell = table[(i, j)]
            cell.set_facecolor('#ecf0f1')
            cell.set_text_props(fontsize=10)

    ax.set_title('Performance Metrics: RMSE (MAE) in µg/m³',
                 fontsize=16, fontweight='bold', pad=30)

    plt.tight_layout()

    img_stream = io.BytesIO()
    plt.savefig(img_stream, format='png', dpi=300, bbox_inches='tight', facecolor='white')
    plt.close()
    img_stream.seek(0)

    return img_stream

def create_presentation():
    """Create the PowerPoint presentation."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
    prs.slide_height = Inches(7.5)

    # Load metrics
    print("📊 Loading metrics...")
    metrics = load_metrics()

    # ==================== SLIDE 1: Title ====================
    print("📝 Creating title slide...")
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Background gradient (approximated with shape)
    background = slide.shapes.add_shape(
        1,  # Rectangle
        0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(102, 126, 234)  # #667eea
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(11.333), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = "TopoFlow"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(72)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(4.2), Inches(11.333), Inches(1)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Physics-Informed Deep Learning for Multi-Pollutant Air Quality Forecasting"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.alignment = PP_ALIGN.CENTER
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = RGBColor(255, 255, 255)

    # Author
    author_box = slide.shapes.add_textbox(
        Inches(1), Inches(5.5), Inches(11.333), Inches(0.5)
    )
    author_frame = author_box.text_frame
    author_frame.text = "Ammar Kheddar | 2025"
    author_para = author_frame.paragraphs[0]
    author_para.alignment = PP_ALIGN.CENTER
    author_para.font.size = Pt(20)
    author_para.font.color.rgb = RGBColor(255, 255, 255)

    # ==================== SLIDE 2: Methodology ====================
    print("🔬 Creating methodology slide...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "Methodology: Physics-Informed Innovations"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(44, 62, 80)

    # Architecture diagram
    print("  → Generating architecture diagram...")
    arch_img = create_architecture_diagram()
    slide.shapes.add_picture(arch_img, Inches(0.5), Inches(1.2), width=Inches(12.333))

    # Innovation boxes (text only)
    innovations = [
        ("🌬️ Wind-Guided Scanning", "Patches reordered by wind direction (upwind→downwind)\n16 pre-computed sectors | Captures atmospheric transport"),
        ("🏔️ Elevation-Based Attention", "Terrain-aware bias: bias = -α × ReLU(Δh/H)\nPenalizes uphill transport | Learnable parameter α"),
        ("🧪 Multi-Pollutant Modeling", "6 species: PM2.5, PM10, SO₂, NO₂, CO, O₃\n4 horizons: 12h, 24h, 48h, 96h"),
        ("⚡ Large-Scale Training", "800 AMD MI250X GPUs on LUMI\n100 nodes × 8 GPUs | ~48 hours")
    ]

    y_start = 4.5
    box_height = 0.65
    spacing = 0.05

    colors_rgb = [
        RGBColor(155, 89, 182),   # Purple
        RGBColor(231, 76, 60),    # Red
        RGBColor(52, 152, 219),   # Blue
        RGBColor(26, 188, 156)    # Teal
    ]

    for i, ((title, desc), color) in enumerate(zip(innovations, colors_rgb)):
        y = y_start + i * (box_height + spacing)

        # Background box
        box_shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0.5), Inches(y), Inches(12.333), Inches(box_height)
        )
        box_shape.fill.solid()
        box_shape.fill.fore_color.rgb = color
        box_shape.line.color.rgb = RGBColor(255, 255, 255)
        box_shape.line.width = Pt(2)

        # Title
        title_textbox = slide.shapes.add_textbox(
            Inches(0.7), Inches(y + 0.05), Inches(12), Inches(0.25)
        )
        tf = title_textbox.text_frame
        tf.text = title
        p = tf.paragraphs[0]
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        # Description
        desc_textbox = slide.shapes.add_textbox(
            Inches(0.7), Inches(y + 0.3), Inches(12), Inches(0.3)
        )
        tf2 = desc_textbox.text_frame
        tf2.text = desc
        p2 = tf2.paragraphs[0]
        p2.font.size = Pt(12)
        p2.font.color.rgb = RGBColor(255, 255, 255)

    # ==================== SLIDE 3: Results - Performance Charts ====================
    print("📊 Creating results slide (charts)...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "Results: Performance Metrics (Test Year 2018)"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(44, 62, 80)

    # RMSE Chart
    print("  → Generating RMSE chart...")
    rmse_img = create_rmse_chart(metrics)
    slide.shapes.add_picture(rmse_img, Inches(0.5), Inches(1.2), width=Inches(12.333))

    # ==================== SLIDE 4: Results - MAE Evolution ====================
    print("📈 Creating MAE evolution slide...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "Mean Absolute Error (MAE) Evolution by Horizon"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(44, 62, 80)

    # MAE Chart
    print("  → Generating MAE chart...")
    mae_img = create_mae_chart(metrics)
    slide.shapes.add_picture(mae_img, Inches(0.5), Inches(1.1), width=Inches(12.333))

    # ==================== SLIDE 5: Results - Comparison Table ====================
    print("📋 Creating comparison table slide...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = "Detailed Performance: RMSE (MAE) Summary"
    title_para = title_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.CENTER
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(44, 62, 80)

    # Table
    print("  → Generating comparison table...")
    table_img = create_comparison_table(metrics)
    slide.shapes.add_picture(table_img, Inches(1), Inches(1.3), width=Inches(11.333))

    # Stats footer
    stats_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12.333), Inches(0.8))
    stats_frame = stats_box.text_frame
    stats_frame.text = "📍 China Region: 11,317 active pixels (34.5%) | 128×256 grid | 4,000 test samples | Val Loss: 0.3557"
    stats_para = stats_frame.paragraphs[0]
    stats_para.alignment = PP_ALIGN.CENTER
    stats_para.font.size = Pt(14)
    stats_para.font.bold = True
    stats_para.font.color.rgb = RGBColor(52, 73, 94)

    # ==================== SLIDE 6: China Maps (if available) ====================
    print("🗺️ Checking for China maps...")
    pm25_map = MEDIA_DIR / "map_pm25_24d.png"
    if pm25_map.exists():
        print("  → Adding China map slide...")
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
        title_frame = title_box.text_frame
        title_frame.text = "Spatial Distribution: PM2.5 Predictions (24h Horizon)"
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(44, 62, 80)

        # Add map
        slide.shapes.add_picture(str(pm25_map), Inches(1.5), Inches(1.3), width=Inches(10.333))

    # Save presentation
    output_file = "TopoFlow_Presentation.pptx"
    prs.save(output_file)
    print(f"\n✅ Presentation created: {output_file}")
    print(f"📊 Total slides: {len(prs.slides)}")

    return output_file

if __name__ == "__main__":
    print("\n" + "="*70)
    print("TOPOFLOW POWERPOINT GENERATOR")
    print("="*70)
    print()

    output = create_presentation()

    print("\n" + "="*70)
    print(f"🎉 SUCCESS! PowerPoint saved to: {output}")
    print("="*70)
    print()
    print("📖 To view:")
    print(f"   1. Copy to your local machine: scp lumi:{Path.cwd()}/{output} .")
    print(f"   2. Open with PowerPoint or LibreOffice Impress")
    print()
